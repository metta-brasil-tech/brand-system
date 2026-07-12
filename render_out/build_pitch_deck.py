#!/usr/bin/env python3
"""Deck de apresentação do Brand System / ad-generator — pra demo de sexta.
Construído do zero (não usa o template metta-deck, que é pra proposta de cliente).
Usa a paleta real (tokens.css) e screenshots/criativos reais gerados hoje.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
FR = ROOT / "render_out" / "demo-video" / "frames"
ADS = ROOT / "assets" / "generated" / "ads"
OUT = ROOT / "render_out" / "brand-system-apresentacao.pptx"

# --- paleta (de styles/tokens.css) ---
YELLOW = RGBColor(0xFF, 0xBE, 0x18)
YELLOW_SOFT = RGBColor(0xFF, 0xE3, 0xA6)
NIGHT_5 = RGBColor(0x0A, 0x10, 0x13)
NIGHT_15 = RGBColor(0x13, 0x1F, 0x25)
NIGHT_20 = RGBColor(0x1A, 0x2A, 0x35)
NIGHT_40 = RGBColor(0x43, 0x59, 0x65)
NIGHT_70 = RGBColor(0x94, 0xB5, 0xC8)
NIGHT_97 = RGBColor(0xEF, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x18, 0x1F)

FONT_DISPLAY = "Arial Black"
FONT_BODY = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def add_text(slide, text, left, top, width, height, size=24, color=INK, bold=False,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, color, line=False):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.shadow.inherit = False
    if line:
        r.line.color.rgb = color; r.line.width = Pt(0.5)
    else:
        r.line.fill.background()
    return r


def add_pill(slide, text, left, top, width, height, bg, fg, size=13):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    r.adjustments[0] = 0.5
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    tf = r.text_frame
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = True; run.font.name = FONT_BODY
    run.font.color.rgb = fg
    return r


def eyebrow(slide, text, left, top, color=YELLOW):
    return add_text(slide, text.upper(), left, top, Inches(8), Inches(0.4),
                     size=14, color=color, bold=True, font=FONT_BODY)


_PNG_CACHE = ROOT / "render_out" / "_pptx_png_cache"
_PNG_CACHE.mkdir(exist_ok=True)


def _as_png(path: Path) -> Path:
    path = Path(path)
    if path.suffix.lower() == ".png":
        return path
    from PIL import Image
    dst = _PNG_CACHE / (path.stem + ".png")
    if not dst.exists():
        Image.open(path).convert("RGB").save(dst, "PNG")
    return dst


def img_fit(slide, path, left, top, max_w, max_h, shadow=True):
    from PIL import Image
    path = _as_png(path)
    im = Image.open(path)
    ratio = im.width / im.height
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w = max_w; h = int(max_w / ratio)
    else:
        h = max_h; w = int(max_h * ratio)
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if shadow:
        pic.shadow.inherit = False
    return pic


def page_number(slide, n, total=9, dark=False):
    color = WHITE if dark else NIGHT_40
    add_text(slide, f"{n:02d} / {total:02d}", SW - Inches(1.6), SH - Inches(0.55),
              Inches(1.3), Inches(0.35), size=11, color=color, align=PP_ALIGN.RIGHT)


# ============================================================ SLIDE 1 — CAPA
s = add_slide(NIGHT_5)
add_rect(s, 0, 0, Inches(0.18), SH, YELLOW)
add_text(s, "METTA · BRAND SYSTEM", Inches(0.9), Inches(0.9), Inches(8), Inches(0.4),
          size=14, color=YELLOW, bold=True)
add_text(s, "Ferramenta de\ngeração de criativos.", Inches(0.85), Inches(2.5), Inches(11.5), Inches(2.6),
          size=54, color=WHITE, bold=True, font=FONT_DISPLAY, line_spacing=1.05)
add_text(s, "Da marca navegável ao anúncio pronto — com direção de arte, geração de imagem\ne QA automatizado, tudo com a identidade Metta.",
          Inches(0.9), Inches(5.0), Inches(9.5), Inches(1.0), size=18, color=NIGHT_70, line_spacing=1.3)
add_text(s, "Apresentação · Julho 2026", Inches(0.9), Inches(6.6), Inches(6), Inches(0.4),
          size=13, color=NIGHT_40)

# ============================================================ SLIDE 2 — O QUE É
s = add_slide(WHITE)
eyebrow(s, "O que é", Inches(0.9), Inches(0.7))
add_text(s, "Um único sistema, três camadas.", Inches(0.85), Inches(1.1), Inches(11), Inches(0.9),
          size=34, color=INK, bold=True, font=FONT_DISPLAY)

cards = [
    ("Marca navegável", "Manifesto, ICP, identidade verbal e visual, direção de arte — tudo documentado e vivo, não um PDF esquecido.", NIGHT_97, INK),
    ("Gerador de criativos", "Pipeline: briefing → direção de arte → geração de imagem → checagem por visão → crítico anti-slop. Duas marcas: Metta e Tiago Alves.", NIGHT_5, WHITE),
    ("Biblioteca viva", "Catálogo canônico (109 peças de referência) + Criativos Gerados — tudo que a ferramenta produz, com veredito do QA documentado.", NIGHT_97, INK),
]
cw = Inches(3.75); gap = Inches(0.3); x0 = Inches(0.9); y0 = Inches(2.4); ch = Inches(4.3)
for i, (title, desc, bg, fg) in enumerate(cards):
    x = x0 + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, bg)
    add_rect(s, x, y0, cw, Inches(0.12), YELLOW)
    add_text(s, title, x + Inches(0.35), y0 + Inches(0.5), cw - Inches(0.7), Inches(0.9),
              size=22, color=fg, bold=True, font=FONT_DISPLAY, line_spacing=1.05)
    add_text(s, desc, x + Inches(0.35), y0 + Inches(1.5), cw - Inches(0.7), Inches(2.6),
              size=15, color=fg, line_spacing=1.35)
page_number(s, 2)

# ============================================================ SLIDE 3 — PIPELINE
s = add_slide(NIGHT_5)
eyebrow(s, "Como funciona", Inches(0.9), Inches(0.7))
add_text(s, "6 etapas, do texto livre ao PNG final.", Inches(0.85), Inches(1.1), Inches(11), Inches(0.9),
          size=34, color=WHITE, bold=True, font=FONT_DISPLAY)

steps = [
    ("01", "Briefing", "Texto livre PT-BR vira briefing estruturado"),
    ("02", "Direção de arte", "Escolhe estilo + composição, referência do banco"),
    ("03", "Geração de imagem", "gpt-image-2 gera a foto, guiado por modos de falha conhecidos"),
    ("04", "Vision-QA", "Uma IA olha o PNG renderizado e julga: passa ou reprova"),
    ("05", "Crítico anti-slop", "Checa 16 padrões de “cara de IA genérica” contra o banco"),
    ("06", "Entrega", "PNG final + HTML, pronto pra publicar"),
]
cw = Inches(3.9); ch = Inches(1.75); gx = Inches(0.25); gy = Inches(0.25)
x0 = Inches(0.9); y0 = Inches(2.35)
for i, (num, title, desc) in enumerate(steps):
    col = i % 3; row = i // 3
    x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
    add_rect(s, x, y, cw, ch, NIGHT_15)
    add_text(s, num, x + Inches(0.3), y + Inches(0.18), Inches(1.2), Inches(0.5),
              size=22, color=YELLOW, bold=True, font=FONT_DISPLAY)
    add_text(s, title, x + Inches(0.3), y + Inches(0.68), cw - Inches(0.6), Inches(0.4),
              size=16, color=WHITE, bold=True)
    add_text(s, desc, x + Inches(0.3), y + Inches(1.08), cw - Inches(0.6), Inches(0.6),
              size=11.5, color=NIGHT_70, line_spacing=1.2)
page_number(s, 3, dark=True)

# ============================================================ SLIDE 4 — SCREENSHOT: VISÃO GERAL
def screenshot_slide(idx, title, subtitle, img_path, total=9):
    s = add_slide(NIGHT_97)
    eyebrow(s, "Na prática", Inches(0.9), Inches(0.55))
    add_text(s, title, Inches(0.85), Inches(0.92), Inches(9), Inches(0.7),
              size=28, color=INK, bold=True, font=FONT_DISPLAY)
    add_text(s, subtitle, Inches(0.9), Inches(1.55), Inches(9.5), Inches(0.5),
              size=14, color=NIGHT_40, line_spacing=1.3)
    add_rect(s, Inches(0.85), Inches(2.15), Inches(11.6), Inches(5.0), WHITE)
    img_fit(s, img_path, Inches(0.95), Inches(2.25), Inches(11.4), Inches(4.8))
    page_number(s, idx, total)
    return s

screenshot_slide(4, "Um hub, não um PDF.",
    "Marca, audiência, identidade verbal/visual, direção de arte — navegável, sempre atualizado.",
    FR / "f_2.png")

# ============================================================ SLIDE 5 — SCREENSHOT: WIZARD
screenshot_slide(5, "Criar começa simples.",
    "Escolhe a marca (Metta ou Tiago) — o resto do wizard adapta tom, estilos e paleta automaticamente.",
    FR / "f_14.png")

# ============================================================ SLIDE 6 — EXEMPLOS GERADOS
s = add_slide(WHITE)
eyebrow(s, "Resultado", Inches(0.9), Inches(0.6))
add_text(s, "Peças reais, geradas hoje.", Inches(0.85), Inches(0.97), Inches(10), Inches(0.7),
          size=30, color=INK, bold=True, font=FONT_DISPLAY)
imgs = [
    ADS / "i-retrato-autoridade-metodo.webp",
    ADS / "light-surreal-escada-escalam-estacionam.webp",
    ADS / "yellow-bloco-convite-mentoria.webp",
    ADS / "tiago-story-hero-foto-real-estrutura.webp",
]
cw = Inches(2.7); ch = Inches(4.85); gap = Inches(0.22)
x0 = Inches(0.9); y0 = Inches(1.85)
for i, p in enumerate(imgs):
    x = x0 + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, NIGHT_97)
    img_fit(s, p, x + Inches(0.06), y0 + Inches(0.06), cw - Inches(0.12), ch - Inches(0.12))
add_text(s, "Todas passaram por checagem de visão (vision-QA) e crítico anti-slop antes de entrar aqui.",
          Inches(0.9), Inches(6.85), Inches(10.5), Inches(0.4), size=12, color=NIGHT_40)
page_number(s, 6)

# ============================================================ SLIDE 7 — SCREENSHOT: BIBLIOTECA
screenshot_slide(7, "Biblioteca de Criativos Gerados.",
    "Toda peça que a ferramenta produz fica registrada — com o veredito do QA, não só a imagem.",
    FR / "f_22.png")

# ============================================================ SLIDE 8 — SCREENSHOT: CATÁLOGO
screenshot_slide(8, "+ Catálogo canônico: 109 referências.",
    "Banco vivo de peças aprovadas — é o que a IA usa como referência de “mesmo designer”.",
    FR / "f_30.png")

# ============================================================ SLIDE 9 — PRÓXIMOS PASSOS
s = add_slide(NIGHT_5)
eyebrow(s, "Próximos passos", Inches(0.9), Inches(0.7))
add_text(s, "O que vem a seguir.", Inches(0.85), Inches(1.1), Inches(10), Inches(0.9),
          size=34, color=WHITE, bold=True, font=FONT_DISPLAY)
roadmap = [
    "Direção de série completa (paletas P1–P10, guardrails C1–C8) — recuperada do backup, aguardando merge",
    "Recompositor full-bleed — resolve o padrão de foto cortando a pessoa no meio",
    "Mais fotos reais do Tiago (bastidor/selfie) para os modelos que exigem foto real",
    "Calibrar o crítico pra comparar só com referências da mesma família de paleta",
]
y = Inches(2.5)
for item in roadmap:
    add_rect(s, Inches(0.9), y + Inches(0.12), Inches(0.14), Inches(0.14), YELLOW)
    add_text(s, item, Inches(1.25), y, Inches(10.5), Inches(0.7), size=16, color=WHITE, line_spacing=1.25)
    y += Inches(0.95)
page_number(s, 9, dark=True)

prs.save(str(OUT))
print("Salvo em:", OUT)
